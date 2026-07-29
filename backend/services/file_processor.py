
import os
import io
import csv
import logging

import pdfplumber
from pdfplumber.utils.exceptions import PdfminerException
from pdfminer.pdfdocument import PDFPasswordIncorrect
import docx
from docx.document import Document as _DocxDocument
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.oxml.ns import qn
from pptx import Presentation

# Configure logger
logger = logging.getLogger(__name__)

# Below this many characters of extracted text, a document is treated as
# "empty" (scanned/image-only or unreadable) - triggers the OCR fallback for
# PDFs and the low-quality warning surfaced to the user.
MIN_MEANINGFUL_CHARS = 100

# Cap OCR to this many pages so a large scanned document can't blow the
# request timeout (OCR is ~1-3s per page).
MAX_OCR_PAGES = 15

# Cap how many rows of a CSV/Excel sheet are turned into text. A huge
# spreadsheet would otherwise become one enormous blob -> far too many RAG
# chunks -> Groq rate-limit/timeout during indexing. Truncation is flagged
# inline so it's visible rather than silent.
MAX_TABLE_ROWS = 500

# Global ceiling on extracted text across ALL formats (~50k tokens). The row
# cap only bounds spreadsheets; a long PDF or Word report can still blow past
# what the free Groq tier can index in one request. Applied centrally in
# extract_text so every path is protected, with the truncation flagged inline.
MAX_EXTRACTED_CHARS = 200_000

# Encodings tried, in order, for text-like uploads (.txt, .csv). Financial
# exports from Excel are frequently cp1252/latin-1 (currency symbols,
# non-breaking spaces); latin-1 maps all 256 byte values so it never raises
# and acts as the final catch-all before a lossy decode.
_TEXT_ENCODINGS = ("utf-8-sig", "cp1252", "latin-1")

# Image formats read via OCR (a photo/scan of a statement or receipt).
IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp", ".gif")


class FileProcessor:
    """
    Handles extraction of text from various file formats.
    Supported: .pdf, .docx, .pptx, .txt, .csv, .xlsx, .xls, and images
    (.png/.jpg/.jpeg/.tiff/.bmp/.webp/.gif via OCR).

    Tables are preserved as pipe-delimited rows rather than flattened into
    jumbled text - important for financial documents, where the numbers in
    balance sheets / income statements are the most valuable content.

    Scanned/image PDFs (no embedded text) and image uploads fall back to OCR
    when the tesseract binary is available.
    """

    @staticmethod
    def extract_text(file_bytes: bytes, filename: str) -> str:
        """
        Extracts text, choosing the parser from the file's actual content
        (magic bytes) and only falling back to the extension. This way a
        mislabeled upload - a PDF saved as .txt, an .xlsx renamed .csv -
        is still read correctly. Output is capped at MAX_EXTRACTED_CHARS.
        """
        ext = FileProcessor._detect_ext(file_bytes, filename)

        try:
            if ext == ".pdf":
                text = FileProcessor._extract_from_pdf(file_bytes)
            elif ext == ".docx":
                text = FileProcessor._extract_from_docx(file_bytes)
            elif ext == ".pptx":
                text = FileProcessor._extract_from_pptx(file_bytes)
            elif ext == ".csv":
                text = FileProcessor._extract_from_csv(file_bytes)
            elif ext in (".xlsx", ".xls"):
                text = FileProcessor._extract_from_excel(file_bytes)
            elif ext == ".txt":
                text = FileProcessor._decode_text(file_bytes)
            elif ext in IMAGE_EXTS:
                text = FileProcessor._extract_from_image(file_bytes)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
        except Exception as e:
            logger.error(f"Error processing {filename}: {e}")
            raise e

        return FileProcessor._cap_text(text)

    @staticmethod
    def _detect_ext(file_bytes: bytes, filename: str) -> str:
        """
        Determine the real file type from its content signature, falling back
        to the filename extension when content is ambiguous (e.g. plain text,
        where .txt vs .csv can't be told apart from bytes).
        """
        declared = os.path.splitext(filename)[1].lower()
        sig = file_bytes[:8]

        if sig[:4] == b"%PDF":
            return ".pdf"
        # OLE2 compound file - legacy Office; .xls is the only one we support.
        if sig == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
            return ".xls"
        # OOXML/zip container: peek inside to tell docx/xlsx/pptx apart.
        if sig[:4] == b"PK\x03\x04":
            return FileProcessor._detect_ooxml(file_bytes) or declared

        # Images (OCR path).
        if sig[:4] == b"\x89PNG":
            return ".png"
        if sig[:3] == b"\xff\xd8\xff":
            return ".jpg"
        if sig[:2] == b"BM":
            return ".bmp"
        if sig[:4] in (b"II*\x00", b"MM\x00*"):
            return ".tiff"
        if sig[:4] == b"GIF8":
            return ".gif"
        if sig[:4] == b"RIFF" and file_bytes[8:12] == b"WEBP":
            return ".webp"

        return declared

    @staticmethod
    def _detect_ooxml(file_bytes: bytes) -> str | None:
        """Identify a zip-based Office file by the parts it contains."""
        import zipfile
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                names = z.namelist()
        except zipfile.BadZipFile:
            return None
        if any(n.startswith("word/") for n in names):
            return ".docx"
        if any(n.startswith("xl/") for n in names):
            return ".xlsx"
        if any(n.startswith("ppt/") for n in names):
            return ".pptx"
        return None

    @staticmethod
    def _cap_text(text: str) -> str:
        """Bound total extracted text so a very large document can't overwhelm
        RAG indexing; flags the truncation inline rather than silently."""
        if text and len(text) > MAX_EXTRACTED_CHARS:
            return (
                text[:MAX_EXTRACTED_CHARS]
                + f"\n\n[... document truncated to the first "
                f"{MAX_EXTRACTED_CHARS} characters ...]"
            )
        return text

    @staticmethod
    def _decode_text(file_bytes: bytes) -> str:
        """
        Decode text-like bytes, trying UTF-8 first and falling back to the
        common Windows/Latin encodings so non-ASCII characters in finance
        exports (£, €, non-breaking spaces) aren't silently dropped.
        """
        for enc in _TEXT_ENCODINGS:
            try:
                return file_bytes.decode(enc)
            except UnicodeDecodeError:
                continue
        # Should be unreachable (latin-1 never raises), but stay safe.
        return file_bytes.decode("utf-8", errors="ignore")

    @staticmethod
    def looks_empty(text: str) -> bool:
        """True if extraction produced almost no text. Used internally to
        decide whether a PDF needs the OCR fallback."""
        return len(text.strip()) < MIN_MEANINGFUL_CHARS

    @staticmethod
    def extraction_insufficient(text: str, filename: str) -> bool:
        """
        True when a file couldn't be usefully read (warn + skip indexing).

        Only PDFs get the "scanned image" threshold - a thin PDF, even after
        the OCR fallback, likely couldn't be read. For text/CSV/Excel/Word/
        PPTX the user supplied the content directly, so any non-empty
        extraction is valid (a short one-line note must not be rejected).
        """
        stripped = text.strip()
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return len(stripped) < MIN_MEANINGFUL_CHARS
        return len(stripped) == 0

    # -----------------------------------------------------------------
    # Shared helper
    # -----------------------------------------------------------------
    @staticmethod
    def _format_table(rows) -> str:
        """
        Render a list of rows (each a list of cells) as pipe-delimited lines
        so row/column structure survives into the extracted text. Empty rows
        and None cells are handled; a fully-empty table yields "".
        """
        lines = []
        for row in rows or []:
            cells = [
                ("" if cell is None else str(cell)).replace("\n", " ").strip()
                for cell in row
            ]
            if any(cells):
                lines.append("| " + " | ".join(cells) + " |")
        return "\n".join(lines)

    # -----------------------------------------------------------------
    # PDF
    # -----------------------------------------------------------------
    @staticmethod
    def _extract_from_pdf(file_bytes: bytes) -> str:
        parts = []
        try:
            pdf = pdfplumber.open(io.BytesIO(file_bytes))
        except (PDFPasswordIncorrect, PdfminerException) as e:
            # pdfplumber wraps pdfminer errors in PdfminerException, so a
            # password failure can arrive either directly or as the wrapped
            # __context__. Map only that case to a clear message; let other
            # parse failures (corrupt PDF, etc.) propagate unchanged.
            cause = e if isinstance(e, PDFPasswordIncorrect) else e.__context__
            if isinstance(cause, PDFPasswordIncorrect):
                raise ValueError(
                    "This PDF is password-protected. Please remove the password "
                    "and upload it again."
                )
            raise
        with pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # Locate tables once so we can both (a) exclude their region
                # from the flat text pass and (b) render them structured.
                tables = page.find_tables()
                bboxes = [t.bbox for t in tables]

                # Flat text with the table regions removed. Without this,
                # extract_text() ALSO flattens the numbers inside each table,
                # so every table value would appear twice (once here, once in
                # the formatted table) - inflating tokens and skewing metric
                # extraction with double-counted figures.
                def _outside_tables(obj):
                    for bx0, btop, bx1, bbottom in bboxes:
                        if (obj["x0"] >= bx0 and obj["x1"] <= bx1
                                and obj["top"] >= btop and obj["bottom"] <= bbottom):
                            return False
                    return True

                page_text = (
                    page.filter(_outside_tables).extract_text()
                    if bboxes else page.extract_text()
                )
                if page_text and page_text.strip():
                    # Page marker gives coarse provenance so retrieved chunks
                    # can be tied back to a page.
                    parts.append(f"--- Page {page_num} ---\n{page_text}")

                # Render each detected table with its structure preserved.
                for t in tables:
                    table = t.extract()
                    # pdfplumber over-detects: chart axes / stray gridlines come
                    # back as tiny 1-cell "tables". Keep only genuine data tables
                    # (>= 2 non-empty rows and >= 2 columns).
                    non_empty = [
                        r for r in table
                        if any(c and str(c).strip() for c in r)
                    ]
                    max_cols = max((len(r) for r in non_empty), default=0)
                    if len(non_empty) < 2 or max_cols < 2:
                        continue

                    formatted = FileProcessor._format_table(table)
                    if formatted:
                        parts.append(f"[Table - page {page_num}]\n{formatted}")

        text = "\n\n".join(parts)

        # Scanned/image PDF: no embedded text was found. Try OCR.
        if FileProcessor.looks_empty(text):
            ocr_text = FileProcessor._ocr_pdf(file_bytes)
            if ocr_text:
                return ocr_text

        return text

    @staticmethod
    def _ocr_pdf(file_bytes: bytes) -> str:
        """
        OCR fallback for scanned/image PDFs. Renders each page to an image
        (pypdfium2, via pdfplumber) and runs tesseract. Returns "" if the
        OCR stack (pytesseract + the tesseract binary) isn't available, so
        callers degrade gracefully instead of crashing.
        """
        try:
            import pytesseract
        except ImportError:
            logger.warning("OCR skipped: pytesseract not installed.")
            return ""

        parts = []
        try:
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page_num, page in enumerate(pdf.pages[:MAX_OCR_PAGES], 1):
                    image = page.to_image(resolution=200).original
                    page_text = pytesseract.image_to_string(image)
                    if page_text and page_text.strip():
                        parts.append(f"--- Page {page_num} (OCR) ---\n{page_text.strip()}")
        except pytesseract.TesseractNotFoundError:
            logger.warning("OCR skipped: tesseract binary not installed on this system.")
            return ""
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
            return ""

        return "\n\n".join(parts)

    # -----------------------------------------------------------------
    # DOCX
    # -----------------------------------------------------------------
    @staticmethod
    def _extract_from_docx(file_bytes: bytes) -> str:
        doc = docx.Document(io.BytesIO(file_bytes))
        parts = []

        # Headers / footers first - financial Word docs often park totals,
        # report dates, and disclaimers here, and doc.paragraphs never sees
        # them. Linked headers/footers inherit and yield nothing, so there's
        # no duplication across sections.
        for section in doc.sections:
            for container in (
                section.first_page_header, section.header, section.even_page_header,
                section.first_page_footer, section.footer, section.even_page_footer,
            ):
                FileProcessor._docx_render_container(container, parts)

        # Main body, in order, recursing into nested tables.
        FileProcessor._docx_render_container(doc, parts)

        # Text boxes (DrawingML/VML) aren't reachable from the block tree.
        textbox_text = FileProcessor._docx_textboxes(doc)
        if textbox_text:
            parts.append(f"[Text box]\n{textbox_text}")

        return "\n".join(p for p in parts if p)

    @staticmethod
    def _iter_block_items(container):
        """Yield Paragraphs and Tables in document order from a Document,
        header/footer, or table cell (python-docx exposes .paragraphs and
        .tables separately, losing their interleaved order)."""
        elm = container.element.body if isinstance(container, _DocxDocument) else container._element
        for child in elm.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, container)
            elif child.tag == qn("w:tbl"):
                yield Table(child, container)

    @staticmethod
    def _docx_render_container(container, parts, depth=0):
        """Append a container's paragraphs and tables (recursing into nested
        tables) to parts, in order."""
        for block in FileProcessor._iter_block_items(container):
            if isinstance(block, Paragraph):
                if block.text.strip():
                    parts.append(block.text)
            else:  # Table
                rows = [[cell.text for cell in row.cells] for row in block.rows]
                formatted = FileProcessor._format_table(rows)
                if formatted:
                    label = "Nested table" if depth else "Table"
                    parts.append(f"[{label}]\n{formatted}")
                # A cell may itself hold tables; render those too.
                for row in block.rows:
                    for cell in row.cells:
                        for nested in cell.tables:
                            FileProcessor._docx_render_container_table(nested, parts, depth + 1)

    @staticmethod
    def _docx_render_container_table(table, parts, depth):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        formatted = FileProcessor._format_table(rows)
        if formatted:
            parts.append(f"[Nested table]\n{formatted}")
        for row in table.rows:
            for cell in row.cells:
                for nested in cell.tables:
                    FileProcessor._docx_render_container_table(nested, parts, depth + 1)

    @staticmethod
    def _docx_textboxes(doc) -> str:
        """Collect text from DrawingML/VML text boxes, which live inside a
        run's drawing and so are skipped by paragraph/table iteration."""
        texts = []
        for txbx in doc.element.body.iter(qn("w:txbxContent")):
            words = [node.text for node in txbx.iter(qn("w:t")) if node.text]
            joined = "".join(words).strip()
            if joined:
                texts.append(joined)
        return "\n".join(texts)

    # -----------------------------------------------------------------
    # Images (OCR)
    # -----------------------------------------------------------------
    @staticmethod
    def _extract_from_image(file_bytes: bytes) -> str:
        """
        OCR an uploaded image (a photo/scan of a statement or receipt).
        Returns "" if the OCR stack (pytesseract + Pillow + the tesseract
        binary) isn't available, so callers degrade gracefully.
        """
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            logger.warning("Image OCR skipped: pytesseract/Pillow not installed.")
            return ""

        try:
            image = Image.open(io.BytesIO(file_bytes))
            return pytesseract.image_to_string(image).strip()
        except pytesseract.TesseractNotFoundError:
            logger.warning("Image OCR skipped: tesseract binary not installed on this system.")
            return ""
        except Exception as e:
            logger.warning(f"Image OCR failed: {e}")
            return ""

    # -----------------------------------------------------------------
    # PPTX
    # -----------------------------------------------------------------
    @staticmethod
    def _extract_from_pptx(file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                # A shape can be a table, a text box, or neither (e.g. image).
                if getattr(shape, "has_table", False):
                    rows = [
                        [cell.text for cell in row.cells]
                        for row in shape.table.rows
                    ]
                    formatted = FileProcessor._format_table(rows)
                    if formatted:
                        parts.append(f"[Table - slide {slide_num}]\n{formatted}")
                elif hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)

    # -----------------------------------------------------------------
    # CSV
    # -----------------------------------------------------------------
    @staticmethod
    def _extract_from_csv(file_bytes: bytes) -> str:
        """
        Parse CSV with the csv module (handles quoted fields/commas correctly)
        and render it as a pipe-delimited table, capped at MAX_TABLE_ROWS.
        """
        text = FileProcessor._decode_text(file_bytes)
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)

        truncated = len(rows) > MAX_TABLE_ROWS
        formatted = FileProcessor._format_table(rows[:MAX_TABLE_ROWS])
        if truncated:
            formatted += f"\n[... truncated to the first {MAX_TABLE_ROWS} rows ...]"
        return formatted

    # -----------------------------------------------------------------
    # Excel (.xlsx / .xls)
    # -----------------------------------------------------------------
    @staticmethod
    def _extract_from_excel(file_bytes: bytes) -> str:
        """
        Read every sheet of a workbook and render each as a pipe-delimited
        table, labelled with the sheet name.
        """
        import pandas as pd

        sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None, header=None)
        parts = []
        for sheet_name, df in sheets.items():
            if df.empty:
                continue
            rows = df.where(df.notna(), "").astype(str).values.tolist()

            truncated = len(rows) > MAX_TABLE_ROWS
            formatted = FileProcessor._format_table(rows[:MAX_TABLE_ROWS])
            if formatted:
                if truncated:
                    formatted += f"\n[... truncated to the first {MAX_TABLE_ROWS} rows ...]"
                parts.append(f"[Sheet: {sheet_name}]\n{formatted}")
        return "\n\n".join(parts)
